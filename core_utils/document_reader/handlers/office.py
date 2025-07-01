from typing import Union, List
import io
from .base import BaseFileHandler
from ..exceptions import MissingDependencyError, FileProcessingError

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


class DocxHandler(BaseFileHandler):
    """Handler for Microsoft Word documents."""

    def can_handle(self, file_extension: str) -> bool:
        return file_extension == ".docx"

    def read_content(
        self, file: Union[str, io.IOBase], return_pages: bool = False
    ) -> Union[str, List[str]]:
        if Document is None:
            raise MissingDependencyError("python-docx is required to read DOCX files")

        try:
            if hasattr(file, "read"):
                file.seek(0)
                doc = Document(file)
            else:
                doc = Document(file)

            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return paragraphs if return_pages else "\n".join(paragraphs)
        except Exception as e:
            raise FileProcessingError(f"Error reading DOCX file: {e}")


class ExcelHandler(BaseFileHandler):
    """Handler for Microsoft Excel files."""

    def can_handle(self, file_extension: str) -> bool:
        return file_extension == ".xlsx"

    def read_content(
        self, file: Union[str, io.IOBase], return_pages: bool = False
    ) -> Union[str, List[str]]:
        if load_workbook is None:
            raise MissingDependencyError("openpyxl is required to read Excel files")

        try:
            if hasattr(file, "read"):
                file.seek(0)
                workbook = load_workbook(filename=file, read_only=True, data_only=True)
            else:
                workbook = load_workbook(filename=file, read_only=True, data_only=True)

            content = []
            for sheet in workbook:
                sheet_content = []
                for row in sheet.iter_rows(values_only=True):
                    row_content = " ".join(
                        str(cell) for cell in row if cell is not None
                    )
                    if row_content.strip():
                        sheet_content.append(row_content)

                if sheet_content:
                    content.append("\n".join(sheet_content))

            return content if return_pages else "\n".join(content)
        except Exception as e:
            raise FileProcessingError(f"Error reading Excel file: {e}")


class PowerPointHandler(BaseFileHandler):
    """Handler for Microsoft PowerPoint files."""

    def can_handle(self, file_extension: str) -> bool:
        return file_extension in [".ppt", ".pptx"]

    def read_content(
        self, file: Union[str, io.IOBase], return_pages: bool = False
    ) -> Union[str, List[str]]:
        if Presentation is None:
            raise MissingDependencyError(
                "python-pptx is required to read PowerPoint files"
            )

        try:
            if hasattr(file, "read"):
                file.seek(0)
                prs = Presentation(file)
            else:
                prs = Presentation(file)

            content = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    content.append("\n".join(slide_text))

            return content if return_pages else "\n".join(content)
        except Exception as e:
            raise FileProcessingError(f"Error reading PowerPoint file: {e}")
